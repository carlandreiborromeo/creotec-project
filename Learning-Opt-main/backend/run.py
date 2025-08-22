from flask import Flask, request, jsonify, send_from_directory, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
from pptx import Presentation
from openpyxl import load_workbook
from datetime import datetime
import mysql.connector
from dotenv import load_dotenv

from app.routes.auth import auth_bp
from app.routes.generate import bp as generate_bp
from app.routes.upload import bp as upload_bp
from app.routes.excel_generate import excel_bp

import os
import uuid
import io, json, re, traceback, tempfile, shutil
import pandas as pd
import requests

# Load environment variables
load_dotenv()

app = Flask(__name__)
CORS(app)

CORS(app, resources={r"/api/*": {"origins": "*"}})

@app.after_request
def expose_headers(resp):
    resp.headers["Access-Control-Expose-Headers"] = "Content-Disposition"
    return resp

# Configuration from environment variables
DB_CONFIG = {
    'host': os.getenv('DB_HOST', 'localhost'),
    'user': os.getenv('DB_USER', 'your_username'),
    'password': os.getenv('DB_PASSWORD', 'your_password'),
    'database': os.getenv('DB_NAME', 'your_database')
}

UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), "uploads", "templates")
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

GENERATED_FOLDER = os.path.join("static", "generated")
os.makedirs(GENERATED_FOLDER, exist_ok=True)

# Blueprints
app.register_blueprint(auth_bp)
app.register_blueprint(generate_bp)
app.register_blueprint(upload_bp)
app.register_blueprint(excel_bp)

DEFAULT_MAPPING = {}
PLACEHOLDER_RE = re.compile(r"\{([^}]+)\}")
recent_downloads = []

# Database connection function
def get_db_connection():
    """Create and return a database connection"""
    try:
        return mysql.connector.connect(**DB_CONFIG)
    except mysql.connector.Error as e:
        app.logger.error(f"Database connection error: {str(e)}")
        raise

def to_number(val):
    """Convert to int/float if numeric, else return original or None."""
    try:
        if val is None or str(val).strip() == "":
            return None
        num = float(val)
        return int(num) if num.is_integer() else num
    except (ValueError, TypeError):
        return val

# Quick ping
@app.route("/api/ping")
def ping():
    return jsonify(ok=True)

@app.route("/")
def home():
    return "Hello, Creo Certificate Backend!"

def format_value(val, fmt=None):
    return "" if val is None else str(val)

def replace_placeholders_in_cell(text, mapping, rowdict):
    if "YEAR LAST ATTENDED" in text.upper():
        context = None
        up = text.upper()
        if "ELEMENTARY" in up:
            context = "ELEMENTARY"
        elif "SECONDARY" in up:
            context = "SECONDARY"
        elif "TERTIARY" in up:
            context = "TERTIARY"
    else:
        context = None

    def repl(m):
        key = m.group(1)
        mp = mapping.get(key, key)
        if isinstance(mp, dict):
            col = mp.get(context) or mp.get("DEFAULT")
        else:
            col = mp
        val = rowdict.get(col, "")
        return format_value(val)

    return PLACEHOLDER_RE.sub(repl, text)

def replace_placeholders_in_worksheet(ws, mapping, rowdict):
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
        for cell in row:
            if isinstance(cell.value, str) and "{" in cell.value and "}" in cell.value:
                cell.value = replace_placeholders_in_cell(cell.value, mapping, rowdict)

def _safe_sheet_title(s: str, used: set) -> str:
    title = (s or "").strip() or "Row"
    for ch in '[]:*?/\\':
        title = title.replace(ch, "-")
    title = title[:31] or "Row"
    orig = title
    i = 2
    while title in used:
        suffix = f" ({i})"
        title = (orig[: 31 - len(suffix)] + suffix) if len(orig) + len(suffix) > 31 else orig + suffix
        i += 1
    used.add(title)
    return title

def _copy_template_sheet_with_fallback(wb, template_ws, new_title):
    try:
        ws_copy = wb.copy_worksheet(template_ws)
        ws_copy.title = new_title
        return ws_copy
    except Exception as e:
        print("[WARN] copy_worksheet failed; falling back to manual copy:", repr(e))
        ws = wb.create_sheet(title=new_title)
        for rng in template_ws.merged_cells.ranges:
            ws.merge_cells(str(rng))
        for r in range(1, template_ws.max_row + 1):
            for c in range(1, template_ws.max_column + 1):
                v = template_ws.cell(row=r, column=c).value
                if v is not None:
                    ws.cell(row=r, column=c, value=v)
        return ws

@app.route('/generate/certificates', methods=['POST'])
def generate_certificates():
    data = request.json
    template_path = data.get("templatePath")
    output_folder = "static/generated"
    os.makedirs(output_folder, exist_ok=True)

    # Get custom filename from request, or fallback
    custom_filename = data.get("filename")
    if custom_filename:
        filename = f"{custom_filename}.pptx"
    else:
        name = data.get("name", "Certificate")
        filename = f"{name.replace(' ', '_')}_Certificate.pptx"

    output_path = os.path.join(output_folder, filename)

    # Load and customize the PPTX
    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{{" in run.text and "}}" in run.text:
                            key = run.text.replace("{{", "").replace("}}", "").strip()
                            run.text = data.get(key, "")

    prs.save(output_path)

    # Return list with one file
    return jsonify({"files": [filename]})

@app.route('/api/generate', methods=['POST'])
def generate_tesda_excel():
    uploaded_file = request.files.get("file")
    if not uploaded_file:
        return jsonify({"error": "No file uploaded"}), 400

    # Save temporarily
    temp_path = os.path.join(UPLOAD_FOLDER, f"temp_{uuid.uuid4().hex}.xlsx")
    uploaded_file.save(temp_path)

    try:
        # Load Excel
        wb = load_workbook(temp_path)
        ws = wb.active

        # Save to generated folder
        now = datetime.now().strftime("%Y%m%d-%H%M%S")
        output_filename = f"tesda_record_{now}.xlsx"
        output_path = os.path.join(GENERATED_FOLDER, output_filename)
        wb.save(output_path)

        # ✅ Track in recent_downloads with full metadata for frontend
        recent_downloads.insert(0, {
            "type": "tesda",
            "filename": output_filename,
            "timestamp": datetime.fromtimestamp(os.path.getmtime(output_path)).strftime("%Y-%m-%d %H:%M:%S"),
            "url": f"/static/generated/{output_filename}"
        })

        return send_file(output_path, as_attachment=True, download_name=output_filename)

    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

def api_generate_certificates():
    import requests

    # 1. Forward request to internal generator
    response = requests.post('http://localhost:5000/generate/certificates', json=request.get_json())
    result = response.json()

    # 2. Handle error if generation failed
    if response.status_code != 200:
        return jsonify({"error": "Failed to generate certificates"}), 500

    # 3. Get list of generated files
    generated_files = result.get("files", [])

    # 4. ✅ Track each generated file in download history
    for fname in generated_files:
        recent_downloads.append({
            "type": "certificate",
            "filename": fname,
        })

    # 5. Return original result to frontend
    return jsonify(result)

@app.route('/api/certificates', methods=['GET'])
def get_certificates():
    try:
        files = [f for f in os.listdir(GENERATED_FOLDER) if f.endswith(".pptx")]
        files.sort(key=lambda x: os.path.getmtime(os.path.join(GENERATED_FOLDER, x)), reverse=True)
        return jsonify(files)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/tesda', methods=['GET'])
def get_tesda_records():
    try:
        files = [f for f in os.listdir(GENERATED_FOLDER) if f.endswith(".xlsx")]
        files.sort(key=lambda x: os.path.getmtime(os.path.join(GENERATED_FOLDER, x)), reverse=True)
        return jsonify(files)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/download-history", methods=["GET"])
def get_download_history():
    
    folder = os.path.join("static", "generated")
    files = [
        f for f in os.listdir(folder)
        if f.endswith(".pptx") or (f.endswith(".xlsx") and "tesda" in f.lower())
        
    ]

    # Sort by last modified time
    files.sort(key=lambda f: os.path.getmtime(os.path.join(folder, f)), reverse=True)

    history = []
    for f in files:
        file_type = "certificate" if f.endswith(".pptx") else "tesda"
        history.append({
            "type": file_type,
            "filename": f,
            "timestamp": datetime.fromtimestamp(os.path.getmtime(os.path.join(folder, f))).strftime("%Y-%m-%d %H:%M"),
            "url": f"/static/generated/{f}"
        })
    return jsonify(history)

@app.route("/api/certificates", methods=["GET"])
def list_certificates():
    folder = os.path.join("static", "generated")
    if not os.path.exists(folder):
        return jsonify([])

    files = [
        f for f in os.listdir(folder)
        if f.endswith(".pptx") and f != "example.pptx"
    ]

    # Sort by last modified time descending
    files.sort(key=lambda f: os.path.getmtime(os.path.join(folder, f)), reverse=True)

    return jsonify(files)

@app.route('/api/tesda')
def list_tesda_files():
    files = [
        f for f in os.listdir("static/generated")
        if f.endswith(".xlsx") and "TESDA" in f
    ]
    return jsonify(files)

# TESDA GENERATION ROUTE (internal)
@app.route('/generate/tesda', methods=['POST'])
def generate_tesda_file():
    try:
        data = request.get_json()
        template_name = data.get("template")
        entries = data.get("data")

        if not template_name or not entries:
            return jsonify({"error": "Missing template or data"}), 400

        template_path = os.path.join("uploads", "templates", template_name)
        if not os.path.exists(template_path):
            return jsonify({"error": "Template not found"}), 404

        base_wb = load_workbook(template_path)
        template_ws = base_wb.active

        used_titles = set()
        for idx, entry in enumerate(entries):
            candidate_name = entry.get("Name", f"Sheet{idx+1}")
            new_title = _safe_sheet_title(candidate_name, used_titles)
            ws_copy = _copy_template_sheet_with_fallback(base_wb, template_ws, new_title)
            replace_placeholders_in_worksheet(ws_copy, {}, entry)

        base_wb.remove(template_ws)

        filename = f"TESDA_{datetime.now().strftime('%Y-%m-%d_%H:%M')}.xlsx"
        output_path = os.path.join("static", "generated", filename)
        base_wb.save(output_path)

        return jsonify({"files": [filename]}), 200

    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500



# API ROUTE THAT CALLS INTERNAL GENERATOR AND TRACKS HISTORY
# ✅ TESDA generation route (calls internal generator and logs to history)
@app.route('/api/generate-tesda', methods=['POST'])
def api_generate_tesda():
    # Forward request to the internal generator endpoint
    response = requests.post('http://localhost:5000/generate/tesda', json=request.get_json())

    if response.status_code != 200:
        return jsonify({"error": "Failed to generate TESDA file"}), 500

    result = response.json()
    generated_files = result.get("files", [])

    for fname in generated_files:
        file_path = os.path.join("static", "generated", fname)
        if os.path.exists(file_path):  # ✅ Only add to history if file exists
            recent_downloads.insert(0, {
                "type": "tesda",
                "filename": fname,
                "timestamp": datetime.fromtimestamp(os.path.getmtime(file_path)).strftime("%Y-%m-%d %H:%M:%S"),
                "url": f"/static/generated/{fname}"
            })

    return jsonify(result)


# ✅ Used by frontend to track downloads and update history
@app.route("/api/download-history", methods=["POST"])
def update_download_history():
    data = request.get_json()
    filename = data.get("filename")
    if not filename:
        return jsonify({"error": "Missing filename"}), 400

    file_path = os.path.join("static", "generated", filename)
    if not os.path.exists(file_path):
        return jsonify({"error": "File does not exist"}), 404

    # Avoid duplicates
    if not any(d.get("filename") == filename for d in recent_downloads):
        file_type = "tesda" if filename.lower().endswith(".xlsx") else "certificate"
        recent_downloads.insert(0, {
            "type": file_type,
            "filename": filename,
            "timestamp": datetime.fromtimestamp(os.path.getmtime(file_path)).strftime("%Y-%m-%d %H:%M"),
            "url": f"/static/generated/{filename}"
        })

    return jsonify({"success": True})


# =================== EXCEL GENERATION WITH DATABASE INTEGRATION ===================

@app.route('/api/generate/excel', methods=['POST'])
def generate_excel():
    temp_dir = None
    connection = None
    try:
        # Validate request
        if not request.is_json:
            return jsonify({"error": "Request must be JSON"}), 400
            
        students = request.json.get("students", [])
        if not students:
            return jsonify({"error": "No student data received"}), 400

        # Load template
        template_path = os.path.join("uploads", "templates", "Grades.xlsx")
        if not os.path.exists(template_path):
            return jsonify({"error": f"Template file not found at {template_path}"}), 500
            
        try:
            wb = load_workbook(template_path)
        except Exception as e:
            return jsonify({"error": f"Failed to load Excel template: {str(e)}"}), 500

        # Get file info from first student
        first_student = students[0] if students else {}
        immersion_date = first_student.get("date_of_immersion", "")
        batch = str(first_student.get("batch", ""))
        school = str(first_student.get("school", ""))

        # Generate Excel file (existing logic)
        sheet_map = {
            "PRODUCTION": wb["PRODUCTION"],
            "SUPPORT": wb["SUPPORT"],
            "TECHNICAL": wb["TECHNICAL"]
        }

        # Fill header cells for all sheets
        for ws in wb.worksheets:
            ws['H8'] = batch + " - " + school
            ws['H9'] = "Date of Immersion: " + immersion_date

        # Fill student data
        row_counter = {"PRODUCTION": 10, "SUPPORT": 10, "TECHNICAL": 10}

        for s in students:
            dept_raw = (s.get("department") or "").strip().upper()
            if dept_raw in ["TECHNICAL", "IT"]:
                dept = "TECHNICAL"
            elif dept_raw in ["PRODUCTION", "PROD"]:
                dept = "PRODUCTION"
            else:
                dept = "SUPPORT"

            ws = sheet_map[dept]
            row = row_counter[dept]

            # Fill basic info
            ws[f'B{row}'] = s.get("last_name", "")
            ws[f'C{row}'] = s.get("first_name", "")
            ws[f'D{row}'] = s.get("middle_name", "")
            ws[f'E{row}'] = s.get("strand", "")
            ws[f'F{row}'] = s.get("department", "")
            ws[f'G{row}'] = to_number(s.get("over_all", ""))

            # Fill grades
            ws[f'H{row}'] = to_number(s.get("WI", ""))
            ws[f'I{row}'] = to_number(s.get("CO", ""))
            ws[f'J{row}'] = to_number(s.get("5S", ""))
            ws[f'K{row}'] = to_number(s.get("BO", ""))
            ws[f'L{row}'] = to_number(s.get("CBO", ""))
            ws[f'M{row}'] = to_number(s.get("SDG", ""))
            ws[f'N{row}'] = to_number(s.get("OHSA", ""))
            ws[f'O{row}'] = to_number(s.get("WE", ""))
            ws[f'P{row}'] = to_number(s.get("UJC", ""))
            ws[f'Q{row}'] = to_number(s.get("ISO", ""))
            ws[f'R{row}'] = to_number(s.get("PO", ""))
            ws[f'S{row}'] = to_number(s.get("HR", ""))
            ws[f'AC{row}'] = to_number(s.get("DS", ""))

            if dept == "PRODUCTION":
                ws[f'V{row}'] = to_number(s.get("WI2", ""))
                ws[f'W{row}'] = to_number(s.get("ELEX", ""))
                ws[f'X{row}'] = to_number(s.get("CM", ""))
                ws[f'Y{row}'] = to_number(s.get("SPC", ""))
                ws[f'AB{row}'] = to_number(s.get("PROD", ""))           

            if dept == "SUPPORT":
                ws[f'U{row}'] = to_number(s.get("PerDev", ""))
                ws[f'Z{row}'] = to_number(s.get("Supp", ""))

            if dept == "TECHNICAL":
                ws[f'T{row}'] = to_number(s.get("AppDev", ""))
                ws[f'AA{row}'] = to_number(s.get("Tech", ""))

            row_counter[dept] += 1

        # Save Excel file
        temp_dir = tempfile.mkdtemp()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_batch = batch.replace(" ", "_") if batch else "Batch"
        filename = f"Immersion_Grades_{safe_batch}_{timestamp}.xlsx"
        output_path = os.path.join(temp_dir, filename)
        wb.save(output_path)

        # Save to database
        try:
            connection = get_db_connection()
            cursor = connection.cursor()

            # Insert into generated_files table
            file_insert_query = """
            INSERT INTO generated_files (filename, original_filename, file_type, batch, school, 
                                       date_of_immersion, total_students, file_path, file_size)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s)
            """
            
            file_size = os.path.getsize(output_path)
            immersion_date_parsed = None
            if immersion_date:
                try:
                    immersion_date_parsed = datetime.strptime(immersion_date, "%Y-%m-%d").date()
                except:
                    pass

            cursor.execute(file_insert_query, (
                filename, filename, 'grades', batch, school,
                immersion_date_parsed, len(students), output_path, file_size
            ))
            
            file_id = cursor.lastrowid

            # Insert students data
            student_insert_query = """
            INSERT INTO generated_file_students (
                file_id, last_name, first_name, middle_name, strand, department,
                school, batch, date_of_immersion, over_all, WI, CO, `5S`, BO, CBO, SDG,
                OHSA, WE, UJC, ISO, PO, HR, DS, WI2, ELEX, CM, SPC, PROD,
                PerDev, Supp, AppDev, Tech
            ) VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s,
                      %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
            """

            for s in students:
                cursor.execute(student_insert_query, (
                    file_id, s.get("last_name", ""), s.get("first_name", ""), s.get("middle_name", ""),
                    s.get("strand", ""), s.get("department", ""), school, batch, immersion_date_parsed,
                    to_number(s.get("over_all")), to_number(s.get("WI")), to_number(s.get("CO")),
                    to_number(s.get("5S")), to_number(s.get("BO")), to_number(s.get("CBO")),
                    to_number(s.get("SDG")), to_number(s.get("OHSA")), to_number(s.get("WE")),
                    to_number(s.get("UJC")), to_number(s.get("ISO")), to_number(s.get("PO")),
                    to_number(s.get("HR")), to_number(s.get("DS")), to_number(s.get("WI2")),
                    to_number(s.get("ELEX")), to_number(s.get("CM")), to_number(s.get("SPC")),
                    to_number(s.get("PROD")), to_number(s.get("PerDev")), to_number(s.get("Supp")),
                    to_number(s.get("AppDev")), to_number(s.get("Tech"))
                ))

            # Log the operation
            log_query = """
            INSERT INTO file_operations_log (file_id, operation_type, operation_details)
            VALUES (%s, %s, %s)
            """
            cursor.execute(log_query, (file_id, 'create', json.dumps({'students_count': len(students)})))

            connection.commit()
            
            # Track in recent downloads
            recent_downloads.insert(0, {
                "type": "grades",
                "filename": filename,
                "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "url": f"/static/generated/{filename}",
                "file_id": file_id
            })

        except Exception as db_error:
            app.logger.error(f"Database error: {str(db_error)}")
            # Continue with file generation even if database fails
            if connection:
                connection.rollback()

        return send_file(
            output_path,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        if connection:
            connection.rollback()
        app.logger.error(f"Error in generate_excel: {str(e)}\n{traceback.format_exc()}")
        return jsonify({"error": "Failed to generate Excel file", "details": str(e)}), 500
        
    finally:
        if connection:
            connection.close()
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except Exception as e:
                app.logger.error(f"Error cleaning up temp files: {str(e)}")

# =================== DATABASE CRUD OPERATIONS ===================

# Get all generated files
@app.route('/api/generated-files', methods=['GET'])
def get_generated_files():
    connection = None
    try:
        connection = get_db_connection()
        cursor = connection.cursor(dictionary=True)
        
        query = """
        SELECT gf.*, 
               COUNT(gfs.id) as student_count,
               AVG(gfs.over_all) as average_performance
        FROM generated_files gf
        LEFT JOIN generated_file_students gfs ON gf.id = gfs.file_id
        WHERE gf.status = 'active'
        GROUP BY gf.id
        ORDER BY gf.created_at DESC
        """
        
        cursor.execute(query)
        files = cursor.fetchall()
        
        # Convert datetime objects to strings for JSON serialization
        for file in files:
            if file['created_at']:
                file['created_at'] = file['created_at'].isoformat()
            if file['updated_at']:
                file['updated_at'] = file['updated_at'].isoformat()
            if file['date_of_immersion']:
                file['date_of_immersion'] = file['date_of_immersion'].isoformat()
                
        return jsonify({"files": files}), 200
        
    except Exception as e:
        app.logger.error(f"Error getting generated files: {str(e)}")
        return jsonify({"error": "Failed to fetch files"}), 500
    finally:
        if connection:
            connection.close()

# Get specific file with student data
@app.route('/api/generated-files/<int:file_id>', methods=['GET'])
def get_file_details(file_id):
    connection = None
    try:
        connection = get_db_connection()
        cursor = connection.cursor(dictionary=True)
        
        # Get file info
        file_query = "SELECT * FROM generated_files WHERE id = %s AND status = 'active'"
        cursor.execute(file_query, (file_id,))
        file_info = cursor.fetchone()
        
        if not file_info:
            return jsonify({"error": "File not found"}), 404
        
        # Get student data
        students_query = "SELECT * FROM generated_file_students WHERE file_id = %s ORDER BY last_name, first_name"
        cursor.execute(students_query, (file_id,))
        students = cursor.fetchall()
        
        # Convert datetime objects to strings
        if file_info['created_at']:
            file_info['created_at'] = file_info['created_at'].isoformat()
        if file_info['updated_at']:
            file_info['updated_at'] = file_info['updated_at'].isoformat()
        if file_info['date_of_immersion']:
            file_info['date_of_immersion'] = file_info['date_of_immersion'].isoformat()
            
        for student in students:
            if student['created_at']:
                student['created_at'] = student['created_at'].isoformat()
            if student['updated_at']:
                student['updated_at'] = student['updated_at'].isoformat()
            if student['date_of_immersion']:
                student['date_of_immersion'] = student['date_of_immersion'].isoformat()
        
        return jsonify({"file": file_info, "students": students}), 200
        
    except Exception as e:
        app.logger.error(f"Error getting file details: {str(e)}")
        return jsonify({"error": "Failed to fetch file details"}), 500
    finally:
        if connection:
            connection.close()

# Update file and student data
# Replace your existing update_file function with this corrected version

@app.route('/api/generated-files/<int:file_id>', methods=['PUT'])
def update_file(file_id):
    connection = None
    try:
        if not request.is_json:
            return jsonify({"error": "Request must be JSON"}), 400
            
        data = request.json
        students = data.get("students", [])
        
        app.logger.info(f"Updating file {file_id} with {len(students)} students")
        
        connection = get_db_connection()
        cursor = connection.cursor()
        
        # Check if file exists
        cursor.execute("SELECT id FROM generated_files WHERE id = %s AND status = 'active'", (file_id,))
        if not cursor.fetchone():
            return jsonify({"error": "File not found"}), 404
        
        # Update file info if provided
        file_update_fields = []
        file_update_values = []
        
        if 'batch' in data and data['batch']:
            file_update_fields.append("batch = %s")
            file_update_values.append(data['batch'])
        if 'school' in data and data['school']:
            file_update_fields.append("school = %s")
            file_update_values.append(data['school'])
        if 'date_of_immersion' in data and data['date_of_immersion']:
            file_update_fields.append("date_of_immersion = %s")
            try:
                immersion_date = datetime.strptime(data['date_of_immersion'], "%Y-%m-%d").date()
                file_update_values.append(immersion_date)
            except ValueError:
                file_update_values.append(None)
            
        if file_update_fields:
            file_update_fields.append("updated_at = CURRENT_TIMESTAMP")
            file_update_query = f"UPDATE generated_files SET {', '.join(file_update_fields)} WHERE id = %s"
            file_update_values.append(file_id)
            cursor.execute(file_update_query, file_update_values)
            app.logger.info(f"Updated file metadata for file_id {file_id}")
        
        # Update students if provided
        if students:
            updated_count = 0
            for student in students:
                try:
                    if 'id' in student and student['id']:  # Update existing student
                        update_query = """
                        UPDATE generated_file_students SET
                            last_name = %s, first_name = %s, middle_name = %s, strand = %s,
                            department = %s, over_all = %s, WI = %s, CO = %s, `5S` = %s,
                            BO = %s, CBO = %s, SDG = %s, OHSA = %s, WE = %s, UJC = %s,
                            ISO = %s, PO = %s, HR = %s, DS = %s, WI2 = %s, ELEX = %s,
                            CM = %s, SPC = %s, PROD = %s, PerDev = %s, Supp = %s,
                            AppDev = %s, Tech = %s, updated_at = CURRENT_TIMESTAMP
                        WHERE id = %s AND file_id = %s
                        """
                        
                        update_values = (
                            student.get("last_name", ""),
                            student.get("first_name", ""),
                            student.get("middle_name", ""),
                            student.get("strand", ""),
                            student.get("department", ""),
                            to_number(student.get("over_all")),
                            to_number(student.get("WI")),
                            to_number(student.get("CO")),
                            to_number(student.get("5S")),  # This might be the issue
                            to_number(student.get("BO")),
                            to_number(student.get("CBO")),
                            to_number(student.get("SDG")),
                            to_number(student.get("OHSA")),
                            to_number(student.get("WE")),
                            to_number(student.get("UJC")),
                            to_number(student.get("ISO")),
                            to_number(student.get("PO")),
                            to_number(student.get("HR")),
                            to_number(student.get("DS")),
                            to_number(student.get("WI2")),
                            to_number(student.get("ELEX")),
                            to_number(student.get("CM")),
                            to_number(student.get("SPC")),
                            to_number(student.get("PROD")),
                            to_number(student.get("PerDev")),
                            to_number(student.get("Supp")),
                            to_number(student.get("AppDev")),
                            to_number(student.get("Tech")),
                            student['id'],
                            file_id
                        )
                        
                        cursor.execute(update_query, update_values)
                        if cursor.rowcount > 0:
                            updated_count += 1
                            app.logger.info(f"Updated student {student.get('first_name', '')} {student.get('last_name', '')}")
                        else:
                            app.logger.warning(f"No student found with id {student['id']} for file {file_id}")
                    
                    else:
                        # Handle case where student doesn't have an ID (shouldn't happen in update, but just in case)
                        app.logger.warning(f"Student without ID found in update request: {student.get('first_name', '')} {student.get('last_name', '')}")
                        
                except Exception as student_error:
                    app.logger.error(f"Error updating individual student {student.get('first_name', '')} {student.get('last_name', '')}: {str(student_error)}")
                    # Continue with other students instead of failing completely
                    continue
            
            app.logger.info(f"Successfully updated {updated_count} out of {len(students)} students")
        
        # Log the operation
        try:
            log_query = """
            INSERT INTO file_operations_log (file_id, operation_type, operation_details)
            VALUES (%s, %s, %s)
            """
            cursor.execute(log_query, (file_id, 'update', json.dumps({
                'updated_students': len(students),
                'file_metadata_updated': len(file_update_fields) > 0
            })))
        except Exception as log_error:
            app.logger.warning(f"Failed to log operation: {str(log_error)}")
            # Don't fail the whole operation for logging issues
        
        connection.commit()
        app.logger.info(f"File {file_id} update completed successfully")
        return jsonify({"message": "File updated successfully"}), 200
        
    except mysql.connector.Error as db_error:
        if connection:
            connection.rollback()
        app.logger.error(f"Database error updating file {file_id}: {str(db_error)}")
        return jsonify({"error": "Database error occurred", "details": str(db_error)}), 500
        
    except Exception as e:
        if connection:
            connection.rollback()
        app.logger.error(f"Error updating file {file_id}: {str(e)}\n{traceback.format_exc()}")
        return jsonify({"error": "Failed to update file", "details": str(e)}), 500
        
    finally:
        if connection:
            connection.close()

# Delete file (soft delete)
@app.route('/api/generated-files/<int:file_id>', methods=['DELETE'])
def delete_file(file_id):
    connection = None
    try:
        connection = get_db_connection()
        cursor = connection.cursor()
        
        # Check if file exists first
        cursor.execute("SELECT id, filename, file_path FROM generated_files WHERE id = %s", (file_id,))
        file_info = cursor.fetchone()
        
        if not file_info:
            return jsonify({"error": "File not found"}), 404
        
        file_path = file_info[2] if file_info[2] else None
        filename = file_info[1]
        
        app.logger.info(f"Starting hard delete for file_id {file_id} ({filename})")
        
        # Log the operation BEFORE deletion (since we're deleting the file record)
        try:
            log_query = """
            INSERT INTO file_operations_log (file_id, operation_type, operation_details)
            VALUES (%s, %s, %s)
            """
            cursor.execute(log_query, (file_id, 'hard_delete', json.dumps({
                'filename': filename,
                'file_path': file_path
            })))
        except Exception as log_error:
            app.logger.warning(f"Failed to log delete operation: {str(log_error)}")
        
        # Step 1: Delete all student records associated with this file
        cursor.execute("DELETE FROM generated_file_students WHERE file_id = %s", (file_id,))
        deleted_students = cursor.rowcount
        app.logger.info(f"Deleted {deleted_students} student records for file_id {file_id}")
        
        # Step 2: Delete the file record itself
        cursor.execute("DELETE FROM generated_files WHERE id = %s", (file_id,))
        deleted_files = cursor.rowcount
        
        if deleted_files == 0:
            # This shouldn't happen since we checked existence above, but just in case
            connection.rollback()
            return jsonify({"error": "File not found during deletion"}), 404
        
        app.logger.info(f"Deleted file record for file_id {file_id}")
        
        # Step 3: Try to delete the actual file from filesystem (if it exists)
        if file_path and os.path.exists(file_path):
            try:
                os.remove(file_path)
                app.logger.info(f"Deleted physical file: {file_path}")
            except Exception as file_delete_error:
                app.logger.warning(f"Failed to delete physical file {file_path}: {str(file_delete_error)}")
                # Don't fail the whole operation if we can't delete the physical file
        
        # Commit all database changes
        connection.commit()
        
        app.logger.info(f"Successfully completed hard delete for file_id {file_id}")
        
        return jsonify({
            "message": "File and all associated data deleted successfully",
            "details": {
                "deleted_students": deleted_students,
                "deleted_files": deleted_files,
                "physical_file_deleted": file_path and os.path.exists(file_path)
            }
        }), 200
        
    except mysql.connector.Error as db_error:
        if connection:
            connection.rollback()
        app.logger.error(f"Database error during hard delete of file {file_id}: {str(db_error)}")
        return jsonify({"error": "Database error occurred during deletion", "details": str(db_error)}), 500
        
    except Exception as e:
        if connection:
            connection.rollback()
        app.logger.error(f"Error during hard delete of file {file_id}: {str(e)}\n{traceback.format_exc()}")
        return jsonify({"error": "Failed to delete file", "details": str(e)}), 500
        
    finally:
        if connection:
            connection.close()


# Optional: Add a separate endpoint for soft delete if you want both options
@app.route('/api/generated-files/<int:file_id>/soft-delete', methods=['DELETE'])
def soft_delete_file(file_id):
    """Soft delete - marks file as deleted but keeps data"""
    connection = None
    try:
        connection = get_db_connection()
        cursor = connection.cursor()
        
        # Soft delete - just mark as deleted
        cursor.execute("UPDATE generated_files SET status = 'deleted', updated_at = CURRENT_TIMESTAMP WHERE id = %s", (file_id,))
        
        if cursor.rowcount == 0:
            return jsonify({"error": "File not found"}), 404
        
        # Log the operation
        log_query = """
        INSERT INTO file_operations_log (file_id, operation_type, operation_details)
        VALUES (%s, %s, %s)
        """
        cursor.execute(log_query, (file_id, 'soft_delete', json.dumps({})))
        
        connection.commit()
        app.logger.info(f"Soft deleted file_id {file_id}")
        
        return jsonify({"message": "File marked as deleted successfully"}), 200
        
    except Exception as e:
        if connection:
            connection.rollback()
        app.logger.error(f"Error soft deleting file: {str(e)}")
        return jsonify({"error": "Failed to delete file"}), 500
    finally:
        if connection:
            connection.close()


# Optional: Add endpoint to permanently delete soft-deleted files
@app.route('/api/generated-files/cleanup-deleted', methods=['DELETE'])
def cleanup_deleted_files():
    """Permanently delete all soft-deleted files and their data"""
    connection = None
    try:
        connection = get_db_connection()
        cursor = connection.cursor(dictionary=True)
        
        # Get all soft-deleted files
        cursor.execute("SELECT id, filename, file_path FROM generated_files WHERE status = 'deleted'")
        deleted_files = cursor.fetchall()
        
        if not deleted_files:
            return jsonify({"message": "No deleted files to cleanup"}), 200
        
        total_students_deleted = 0
        total_files_deleted = 0
        physical_files_deleted = 0
        
        for file_info in deleted_files:
            file_id = file_info['id']
            file_path = file_info['file_path']
            
            # Delete student records
            cursor.execute("DELETE FROM generated_file_students WHERE file_id = %s", (file_id,))
            total_students_deleted += cursor.rowcount
            
            # Delete file record
            cursor.execute("DELETE FROM generated_files WHERE id = %s", (file_id,))
            total_files_deleted += cursor.rowcount
            
            # Delete physical file
            if file_path and os.path.exists(file_path):
                try:
                    os.remove(file_path)
                    physical_files_deleted += 1
                except Exception as e:
                    app.logger.warning(f"Failed to delete physical file {file_path}: {str(e)}")
        
        connection.commit()
        
        return jsonify({
            "message": "Cleanup completed successfully",
            "details": {
                "files_processed": len(deleted_files),
                "students_deleted": total_students_deleted,
                "files_deleted": total_files_deleted,
                "physical_files_deleted": physical_files_deleted
            }
        }), 200
        
    except Exception as e:
        if connection:
            connection.rollback()
        app.logger.error(f"Error during cleanup: {str(e)}")
        return jsonify({"error": "Failed to cleanup deleted files"}), 500
    finally:
        if connection:
            connection.close()

# Regenerate and download Excel file
@app.route('/api/generated-files/<int:file_id>/download', methods=['GET'])
def download_file(file_id):
    connection = None
    temp_dir = None
    try:
        connection = get_db_connection()
        cursor = connection.cursor(dictionary=True)
        
        # Get file and student data
        file_query = "SELECT * FROM generated_files WHERE id = %s AND status = 'active'"
        cursor.execute(file_query, (file_id,))
        file_info = cursor.fetchone()
        
        if not file_info:
            return jsonify({"error": "File not found"}), 404
        
        students_query = "SELECT * FROM generated_file_students WHERE file_id = %s"
        cursor.execute(students_query, (file_id,))
        students = cursor.fetchall()
        
        # Regenerate Excel file with current data
        template_path = os.path.join("uploads", "templates", "Grades.xlsx")
        if not os.path.exists(template_path):
            return jsonify({"error": "Template file not found"}), 500
            
        wb = load_workbook(template_path)
        
        # Map departments to worksheets
        sheet_map = {
            "PRODUCTION": wb["PRODUCTION"],
            "SUPPORT": wb["SUPPORT"],
            "TECHNICAL": wb["TECHNICAL"]
        }

        # Fill header cells
        batch = file_info.get('batch', '')
        school = file_info.get('school', '')
        date_of_immersion = file_info.get('date_of_immersion', '')
        if date_of_immersion:
            date_of_immersion = date_of_immersion.strftime("%Y-%m-%d") if hasattr(date_of_immersion, 'strftime') else str(date_of_immersion)
        
        for ws in wb.worksheets:
            ws['H8'] = f"{batch} - {school}"
            ws['H9'] = f"Date of Immersion: {date_of_immersion}"

        # Fill student data
        row_counter = {"PRODUCTION": 10, "SUPPORT": 10, "TECHNICAL": 10}

        for s in students:
            dept_raw = (s.get("department") or "").strip().upper()
            if dept_raw in ["TECHNICAL", "IT"]:
                dept = "TECHNICAL"
            elif dept_raw in ["PRODUCTION", "PROD"]:
                dept = "PRODUCTION"
            else:
                dept = "SUPPORT"

            ws = sheet_map[dept]
            row = row_counter[dept]

            # Fill all the data (same as generate_excel)
            ws[f'B{row}'] = s.get("last_name", "")
            ws[f'C{row}'] = s.get("first_name", "")
            ws[f'D{row}'] = s.get("middle_name", "")
            ws[f'E{row}'] = s.get("strand", "")
            ws[f'F{row}'] = s.get("department", "")
            ws[f'G{row}'] = to_number(s.get("over_all", ""))
            ws[f'H{row}'] = to_number(s.get("WI", ""))
            ws[f'I{row}'] = to_number(s.get("CO", ""))
            ws[f'J{row}'] = to_number(s.get("5S", ""))
            ws[f'K{row}'] = to_number(s.get("BO", ""))
            ws[f'L{row}'] = to_number(s.get("CBO", ""))
            ws[f'M{row}'] = to_number(s.get("SDG", ""))
            ws[f'N{row}'] = to_number(s.get("OHSA", ""))
            ws[f'O{row}'] = to_number(s.get("WE", ""))
            ws[f'P{row}'] = to_number(s.get("UJC", ""))
            ws[f'Q{row}'] = to_number(s.get("ISO", ""))
            ws[f'R{row}'] = to_number(s.get("PO", ""))
            ws[f'S{row}'] = to_number(s.get("HR", ""))
            ws[f'AC{row}'] = to_number(s.get("DS", ""))

            if dept == "PRODUCTION":
                ws[f'V{row}'] = to_number(s.get("WI2", ""))
                ws[f'W{row}'] = to_number(s.get("ELEX", ""))
                ws[f'X{row}'] = to_number(s.get("CM", ""))
                ws[f'Y{row}'] = to_number(s.get("SPC", ""))
                ws[f'AB{row}'] = to_number(s.get("PROD", ""))

            if dept == "SUPPORT":
                ws[f'U{row}'] = to_number(s.get("PerDev", ""))
                ws[f'Z{row}'] = to_number(s.get("Supp", ""))

            if dept == "TECHNICAL":
                ws[f'T{row}'] = to_number(s.get("AppDev", ""))
                ws[f'AA{row}'] = to_number(s.get("Tech", ""))

            row_counter[dept] += 1
        
        # Save temporary file and send
        temp_dir = tempfile.mkdtemp()
        output_path = os.path.join(temp_dir, file_info['filename'])
        wb.save(output_path)
        
        # Log download
        log_query = """
        INSERT INTO file_operations_log (file_id, operation_type, operation_details)
        VALUES (%s, %s, %s)
        """
        cursor.execute(log_query, (file_id, 'download', json.dumps({})))
        connection.commit()
        
        return send_file(
            output_path,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name=file_info['filename']
        )
        
    except Exception as e:
        app.logger.error(f"Error downloading file: {str(e)}")
        return jsonify({"error": "Failed to download file", "details": str(e)}), 500
    finally:
        if connection:
            connection.close()
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except Exception as e:
                app.logger.error(f"Error cleaning up temp files: {str(e)}")
                # Add this to the very end of your run.py file

if __name__ == '__main__':
    # Configuration
    port = int(os.environ.get('PORT', 5000))
    debug = os.environ.get('DEBUG', 'True').lower() == 'true'
    host = os.environ.get('HOST', '127.0.0.1')
    
    print("="*60)
    print("🚀 CREO Certificate Backend Starting...")
    print("="*60)
    print(f"📡 Server: http://{host}:{port}")
    print(f"🐛 Debug mode: {debug}")
    print(f"🗄️  Database: {DB_CONFIG['host']}:{DB_CONFIG.get('port', 3306)}")
    print(f"📁 Upload folder: {UPLOAD_FOLDER}")
    print(f"📁 Generated folder: {GENERATED_FOLDER}")
    print()
    print("📊 Available API endpoints:")
    print("  - GET    /                           - Home page")
    print("  - GET    /api/ping                   - Health check")
    print("  - POST   /api/generate/excel         - Generate Excel grades")
    print("  - GET    /api/generated-files        - List all files")
    print("  - GET    /api/generated-files/<id>   - Get file details")
    print("  - PUT    /api/generated-files/<id>   - Update file")
    print("  - DELETE /api/generated-files/<id>   - Delete file")
    print("  - GET    /api/generated-files/<id>/download - Download file")
    print("="*60)
    
    try:
        app.run(
            host=host,
            port=port,
            debug=debug,
            threaded=True,
            use_reloader=False  # Prevents double startup in debug mode
        )
    except KeyboardInterrupt:
        print("\n👋 Server stopped by user")
    except Exception as e:
        print(f"❌ Error starting server: {e}")
        print("💡 Common fixes:")
        print("   - Check if port is already in use: netstat -an | findstr :5000")
        print("   - Try a different port: set PORT=8000 && python run.py")
        print("   - Check firewall settings")
        print("   - Verify all dependencies are installed: pip install -r requirements.txt")